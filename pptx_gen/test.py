import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import matplotlib.pyplot as plt
import numpy as np

# Create a new presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
title_and_content_layout = prs.slide_layouts[1]  # Title and Content layout
section_header_layout = prs.slide_layouts[2]  # Section Header layout
two_content_layout = prs.slide_layouts[3]  # Two Content layout
blank_slide_layout = prs.slide_layouts[6]  # Blank Slide layout


# Function to add a picture to a slide
def add_picture(slide, image_path, left, top, width=None, height=None):
    if os.path.isfile(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        # Placeholder rectangle if image not found
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left,
            top,
            width or Inches(2),
            height or Inches(2),
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 200, 200)
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        text_frame = shape.text_frame
        text_frame.text = "Image Placeholder"


# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = (
    "Empowering Innovation: Robotics, IoT, and Their Role in Sustainable Development"
)
subtitle.text = "IEEE SPAC 2024 | North South University, Bangladesh\nRedwan Ferdous, CEO of JRC Board & Director of FronTech LTD"

# Slide 2: Introduction
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Introduction"

content = [
    "Session Objectives:",
    "• Inspire innovation among students.",
    "• Discuss the role of Robotics and IoT in sustainable development.",
    "• Highlight opportunities in Bangladesh with a global perspective.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
p.text = content[0]
p.font.bold = True
for line in content[1:]:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 1  # Indent for bullet points

# Add image placeholder on the right
add_picture(
    slide,
    "students_working.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 3: About the Speaker
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "About Me - Redwan Ferdous"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "CEO of JRC Board.",
    "Director of FronTech LTD.",
    "Experience in leading IoT startups in Bangladesh.",
    "Passionate about innovation and sustainability.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.level = 0

# Right Content (Add your photo)
add_picture(
    slide,
    "redwan_ferdous.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(4),
)

# Slide 4: Agenda
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Agenda"

agenda_items = [
    "1. The Rise of Robotics and IoT.",
    "2. Sustainable Development Goals (SDGs).",
    "3. Intersection of Robotics, IoT, and SDGs.",
    "4. Case Studies in Bangladesh.",
    "5. Global Exposure and Opportunities.",
    "6. Q&A Session.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for item in agenda_items:
    p = text_frame.add_paragraph()
    p.text = item
    p.level = 0

# Slide 5: The Rise of Robotics and IoT
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "The Rise of Robotics and IoT"

content = [
    "• Evolution over the years.",
    "• Key global advancements.",
    "• Impact on various industries.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add a timeline graphic (Placeholder)
add_picture(
    slide,
    "timeline_graphic.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 6: Global Statistics (Graph using matplotlib)
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Global Impact of Robotics and IoT"

# Generate a sample bar chart using matplotlib
years = [2020, 2021, 2022, 2023, 2024, 2025]
market_size = [100, 150, 200, 300, 450, 600]  # Sample data

plt.figure(figsize=(6, 4))
plt.bar(years, market_size, color="skyblue")
plt.title("IoT Market Size Projections")
plt.xlabel("Year")
plt.ylabel("Market Size (in billions USD)")
plt.tight_layout()
graph_path = "iot_market_projection.png"
plt.savefig(graph_path)
plt.close()

# Add graph to slide
add_picture(
    slide, graph_path, left=Inches(2), top=Inches(2), width=Inches(6), height=Inches(4)
)

# Slide 7: Understanding Sustainable Development
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Understanding Sustainable Development"

content = [
    "• Definition of Sustainable Development.",
    "• Introduction to the UN Sustainable Development Goals (SDGs).",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add SDG wheel image (Placeholder)
add_picture(
    slide,
    "sdg_wheel.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(4),
)

# Slide 8: SDGs Relevant to Technology
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "SDGs and Technology"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
sdgs = [
    "SDG 7: Affordable and Clean Energy",
    "SDG 9: Industry, Innovation, and Infrastructure",
    "SDG 11: Sustainable Cities and Communities",
]
for sdg in sdgs:
    p = left_text_frame.add_paragraph()
    p.text = sdg
    p.level = 0

# Right Content (Add icons or images for each SDG)
add_picture(
    slide,
    "sdg_icons.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 9: Intersection of Robotics, IoT, and SDGs
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Connecting Robotics and IoT with SDGs"

# Add Venn diagram image (Placeholder)
add_picture(
    slide,
    "venn_diagram.jpg",
    left=Inches(2),
    top=Inches(2),
    width=Inches(6),
    height=Inches(4),
)

# Slide 10: Case Study - Smart Agriculture
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Case Study: Smart Agriculture in Bangladesh"

content = [
    "• Use of IoT sensors for soil monitoring.",
    "• Precision farming techniques.",
    "• Benefits: Increased yields, reduced waste.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add an image of a farmer using IoT devices (Placeholder)
add_picture(
    slide,
    "smart_agriculture.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 11: Graph - Agricultural Productivity
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Impact on Agricultural Productivity"

# Sample data for the graph
categories = ["Before IoT", "After IoT"]
productivity = [2.5, 4.0]  # Sample yield in tons per hectare

plt.figure(figsize=(4, 4))
plt.bar(categories, productivity, color=["orange", "green"])
plt.title("Yield Improvements After IoT Implementation")
plt.ylabel("Yield (tons per hectare)")
plt.tight_layout()
graph_path = "agriculture_productivity.png"
plt.savefig(graph_path)
plt.close()

# Add graph to slide
add_picture(
    slide, graph_path, left=Inches(3), top=Inches(2), width=Inches(6), height=Inches(4)
)

# Slide 12: Case Study - Healthcare
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Case Study: IoT in Healthcare"

content = [
    "• Remote patient monitoring systems.",
    "• Telemedicine for rural communities.",
    "• Outcomes: Improved patient care, accessibility.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add image of a patient using a wearable health device (Placeholder)
add_picture(
    slide,
    "healthcare_iot.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 13: Chart - Healthcare Access
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Improved Healthcare Access"

# Sample data for pie chart
labels = "Before IoT", "After IoT"
sizes = [60, 90]  # Percentage of access

plt.figure(figsize=(4, 4))
plt.pie(
    sizes,
    labels=labels,
    autopct="%1.1f%%",
    startangle=90,
    colors=["lightcoral", "lightskyblue"],
)
plt.title("Healthcare Access Increase")
plt.tight_layout()
graph_path = "healthcare_access.png"
plt.savefig(graph_path)
plt.close()

# Add chart to slide
add_picture(
    slide, graph_path, left=Inches(3), top=Inches(2), width=Inches(6), height=Inches(4)
)

# Slide 14: Role of JRC Board
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Our Journey at JRC Board"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
points = ["• Key projects and innovations.", "• Collaborations with local communities."]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.level = 0

# Right Content (Add images of team and products)
add_picture(
    slide,
    "jrc_board_projects.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 15: Role of FronTech LTD
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Innovations at FronTech LTD"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
points = ["• IoT solutions developed for sustainability.", "• Impact stories."]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.level = 0

# Right Content (Add images of products and solutions)
add_picture(
    slide,
    "frontech_innovations.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 16: Encouraging Innovation
slide = prs.slides.add_slide(blank_slide_layout)
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = "Encouraging Innovation Among Students"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Add background image
add_picture(
    slide,
    "students_innovation.jpg",
    left=Inches(0),
    top=Inches(1.5),
    width=Inches(10),
    height=Inches(5),
)

# Slide 17: Global Exposure
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Global Opportunities"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "• Collaborations with international organizations.",
    "• Participating in global tech challenges.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.level = 0

# Right Content (World map highlighting countries)
add_picture(
    slide,
    "world_map.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 18: Challenges and Solutions
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Challenges in Implementing IoT Solutions"

# Left Column: Challenges
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
p = left_text_frame.paragraphs[0]
p.text = "Challenges:"
p.font.bold = True
challenges = [
    "• Infrastructure limitations.",
    "• High costs.",
    "• Need for skilled personnel.",
]
for challenge in challenges:
    p = left_text_frame.add_paragraph()
    p.text = challenge
    p.level = 1

# Right Column: Solutions
right_shape = slide.shapes.placeholders[2]
right_text_frame = right_shape.text_frame
right_text_frame.clear()
p = right_text_frame.paragraphs[0]
p.text = "Solutions:"
p.font.bold = True
solutions = [
    "• Government support.",
    "• Local manufacturing.",
    "• Education and training.",
]
for solution in solutions:
    p = right_text_frame.add_paragraph()
    p.text = solution
    p.level = 1

# Slide 19: Education and Skill Development
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Education and Skill Development"

content = [
    "• Importance of STEM education.",
    "• Training programs and workshops offered.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add image of a workshop (Placeholder)
add_picture(
    slide,
    "workshop.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 20: Collaborations and Partnerships
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Building Collaborations"

# Left Content
left_text_frame = slide.shapes.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "• Partnerships with universities and industry.",
    "• Benefits of collaborative innovation.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.level = 0

# Right Content (Logos of partners)
add_picture(
    slide,
    "partner_logos.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 21: Future Outlook
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "The Future of Robotics and IoT in Bangladesh"

content = [
    "• Emerging technologies (AI integration, 5G).",
    "• Vision for the next decade.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 0

# Add image of a futuristic smart city (Placeholder)
add_picture(
    slide,
    "future_smart_city.jpg",
    left=Inches(5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 22: How You Can Get Involved
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "How You Can Get Involved"

content = [
    "Opportunities:",
    "• Internships.",
    "• Research projects.",
    "• Innovation challenges.",
    "Steps to participate.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
p.text = content[0]
p.font.bold = True
for line in content[1:4]:
    p = text_frame.add_paragraph()
    p.text = line
    p.level = 1
p = text_frame.add_paragraph()
p.text = content[4]
p.level = 0

# Add image of enthusiastic students (Placeholder)
add_picture(
    slide,
    "students_get_involved.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 23: Call to Action
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Join the Movement"

# Add inspirational quote
text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
quote = '"Innovation is the ability to see change as an opportunity - not a threat." - Steve Jobs'
p = text_frame.paragraphs[0]
p.text = quote
p.alignment = PP_ALIGN.CENTER
p.font.italic = True
p.font.size = Pt(24)

# Add background image (Placeholder)
add_picture(
    slide,
    "call_to_action.jpg",
    left=Inches(0),
    top=Inches(2),
    width=Inches(10),
    height=Inches(4),
)

# Slide 24: Q&A Session
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Questions & Answers"
subtitle = slide.placeholders[1]
subtitle.text = ""

# Add background image (Placeholder)
add_picture(
    slide,
    "q_and_a.jpg",
    left=Inches(0),
    top=Inches(1),
    width=Inches(10),
    height=Inches(5),
)

# Slide 25: Thank You
slide = prs.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Thank You"

content = [
    "Contact Information:",
    "Email: your.email@example.com",
    "LinkedIn: linkedin.com/in/yourprofile",
    "Website: www.yourcompanywebsite.com",
    "Invitation to connect and collaborate.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    if line == content[0]:
        p.font.bold = True

# Add company logos or professional photo (Placeholder)
add_picture(
    slide, "logos.jpg", left=Inches(2), top=Inches(4), width=Inches(6), height=Inches(1)
)

# Save the presentation
prs.save("Empowering_Innovation.pptx")
