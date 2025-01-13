import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np

# Create a new presentation
prs = Presentation()

# Define custom theme colors (optional)
TITLE_COLOR = RGBColor(0, 112, 192)  # A shade of blue for titles
CONTENT_COLOR = RGBColor(60, 60, 60)  # Dark gray for content text
SUBTITLE_COLOR = RGBColor(80, 80, 80)  # Gray for subtitles
BACKGROUND_COLOR = RGBColor(245, 245, 245)  # Light gray background

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
title_and_content_layout = prs.slide_layouts[1]  # Title and Content layout
section_header_layout = prs.slide_layouts[2]  # Section Header layout
two_content_layout = prs.slide_layouts[3]  # Two Content layout
blank_slide_layout = prs.slide_layouts[6]  # Blank Slide layout


# Function to add a picture to a slide
def add_picture(slide, image_path, left, top, width=None, height=None):
    if os.path.isfile(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        # Placeholder rectangle if image not found
        placeholder_width = width or Cm(5)
        placeholder_height = height or Cm(5)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, placeholder_width, placeholder_height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(230, 230, 230)
        line = shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Image Placeholder"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE


# Function to add a title with custom font and size
def add_custom_title(slide, text):
    # Try to access the title placeholder
    title_placeholder = None
    for shape in slide.shapes:
        if (
            shape.is_placeholder and shape.placeholder_format.type == 1
        ):  # Placeholder for title has type 1
            title_placeholder = shape
            break

    if title_placeholder is None:
        # Add a new textbox for the title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(1)
        title_placeholder = slide.shapes.add_textbox(left, top, width, height)

    # Clear any existing text
    title_placeholder.text = ""
    p = title_placeholder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


# Function to add content text with custom styles
def add_custom_content(text_frame, content_lines):
    text_frame.clear()
    for line in content_lines:
        p = text_frame.add_paragraph()
        if line.endswith(":"):
            # This is a heading within content
            run = p.add_run()
            run.text = line
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = TITLE_COLOR
        else:
            # Regular bullet point
            p.text = line
            p.level = 1
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.font.color.rgb = CONTENT_COLOR


# Slide 1: Title Slide with Enhanced Styling
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = (
    "Empowering Innovation: Robotics, IoT, and Their Role in Sustainable Development"
)
subtitle.text = "IEEE SPAC 2024 | North South University, Bangladesh\nRedwan Ferdous, CEO of JRC Board & Director of FronTech LTD"

# Customize title font
title_tf = title.text_frame
title_p = title_tf.paragraphs[0]
title_p.font.name = "Calibri"
title_p.font.size = Pt(48)
title_p.font.bold = True
title_p.font.color.rgb = TITLE_COLOR

# Customize subtitle font
subtitle_tf = subtitle.text_frame
for paragraph in subtitle_tf.paragraphs:
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(24)
        run.font.color.rgb = SUBTITLE_COLOR

# Add background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BACKGROUND_COLOR

# Slide 2: Introduction
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Introduction")

content = [
    "Session Objectives:",
    "Inspire innovation among students.",
    "Discuss the role of Robotics and IoT in sustainable development.",
    "Highlight opportunities in Bangladesh with a global perspective.",
]

text_frame = slide.shapes.placeholders[1].text_frame
add_custom_content(text_frame, content)

# Add image placeholder on the right
add_picture(
    slide,
    "students_working.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 3: About the Speaker
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "About Me - Redwan Ferdous")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "CEO of JRC Board.",
    "Director of FronTech LTD.",
    "Experience in leading IoT startups in Bangladesh.",
    "Passionate about innovation and sustainability.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (Add your photo)
add_picture(
    slide,
    "redwan_ferdous.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(3),
    height=Inches(4),
)

# Slide 4: Agenda
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Agenda")

agenda_items = [
    "1. The Rise of Robotics and IoT.",
    "2. Sustainable Development Goals (SDGs).",
    "3. Intersection of Robotics, IoT, and SDGs.",
    "4. Case Studies in Bangladesh.",
    "5. Global Exposure and Opportunities.",
    "6. Q&A Session.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for item in agenda_items:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Slide 5: The Rise of Robotics and IoT
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "The Rise of Robotics and IoT")

content = [
    "• Evolution over the years.",
    "• Key global advancements.",
    "• Impact on various industries.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR
    p.level = 0

# Add a timeline graphic (Placeholder)
add_picture(
    slide,
    "timeline_graphic.jpg",
    left=Inches(5.5),
    top=Inches(2),
    width=Inches(4),
    height=Inches(3),
)

# Slide 6: Global Statistics (Graph using matplotlib)
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Global Impact of Robotics and IoT")

# Generate a sample bar chart using matplotlib
years = [2020, 2021, 2022, 2023, 2024, 2025]
market_size = [100, 150, 200, 300, 450, 600]  # Sample data

plt.figure(figsize=(6, 4))
plt.bar(years, market_size, color="skyblue")
plt.title("IoT Market Size Projections", fontsize=16)
plt.xlabel("Year", fontsize=12)
plt.ylabel("Market Size (in billions USD)", fontsize=12)
plt.xticks(fontsize=10)
plt.yticks(fontsize=10)
plt.tight_layout()
graph_path = "iot_market_projection.png"
plt.savefig(graph_path, dpi=100)
plt.close()

# Add graph to slide
add_picture(
    slide, graph_path, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)
)

# Slide 7: Understanding Sustainable Development
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Understanding Sustainable Development")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
content = [
    "Definition of Sustainable Development.",
    "Introduction to the UN Sustainable Development Goals (SDGs).",
]
for line in content:
    p = left_text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (SDG wheel image)
add_picture(
    slide,
    "sdg_wheel.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(3),
    height=Inches(3),
)

# Slide 8: SDGs Relevant to Technology
slide = prs.slides.add_slide(blank_slide_layout)

# Add custom title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_tf = title_shape.text_frame
title_p = title_tf.paragraphs[0]
run = title_p.add_run()
run.text = "SDGs and Technology"
run.font.name = "Calibri"
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = TITLE_COLOR

# Add three columns with SDG information
sdg_texts = [
    "SDG 7\nAffordable and Clean Energy",
    "SDG 9\nIndustry, Innovation, and Infrastructure",
    "SDG 11\nSustainable Cities and Communities",
]

for i, sdg_text in enumerate(sdg_texts):
    left = Inches(1 + i * 3)
    top = Inches(2)
    width = Inches(2.5)
    height = Inches(2.5)
    txBox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = txBox.fill
    fill.solid()

    # Assign different colors to each SDG box
    sdg_colors = [RGBColor(74, 134, 232), RGBColor(255, 192, 0), RGBColor(112, 173, 71)]
    fill.fore_color.rgb = sdg_colors[i]

    line = txBox.line
    line.color.rgb = RGBColor(255, 255, 255)

    # Add text
    tf = txBox.text_frame
    tf.text = sdg_text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Slide 9: Intersection of Robotics, IoT, and SDGs
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Connecting Robotics and IoT with SDGs")

# Add Venn diagram image (Placeholder)
add_picture(
    slide,
    "venn_diagram.jpg",
    left=Inches(2),
    top=Inches(2),
    width=Inches(6),
    height=Inches(4),
)

# Slide 10: Case Study - Smart Agriculture
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Case Study: Smart Agriculture in Bangladesh")

content = [
    "• Use of IoT sensors for soil monitoring.",
    "• Precision farming techniques.",
    "• Benefits: Increased yields, reduced waste.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Add an image of a farmer using IoT devices (Placeholder)
add_picture(
    slide,
    "smart_agriculture.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 11: Graph - Agricultural Productivity
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Impact on Agricultural Productivity")

# Sample data for the graph
categories = ["Before IoT", "After IoT"]
productivity = [2.5, 4.0]  # Sample yield in tons per hectare

plt.figure(figsize=(4, 4))
plt.bar(categories, productivity, color=["#ff9999", "#66b3ff"])
plt.title("Yield Improvements After IoT Implementation", fontsize=14)
plt.ylabel("Yield (tons per hectare)", fontsize=12)
plt.tight_layout()
graph_path = "agriculture_productivity.png"
plt.savefig(graph_path, dpi=100)
plt.close()

# Add graph to slide
add_picture(
    slide, graph_path, left=Inches(3), top=Inches(2), width=Inches(5), height=Inches(4)
)

# Slide 12: Case Study - Healthcare
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Case Study: IoT in Healthcare")

content = [
    "• Remote patient monitoring systems.",
    "• Telemedicine for rural communities.",
    "• Outcomes: Improved patient care, accessibility.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Add image of a patient using a wearable health device (Placeholder)
add_picture(
    slide,
    "healthcare_iot.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 13: Chart - Healthcare Access
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Improved Healthcare Access")

# Sample data for pie chart
labels = "Before IoT", "After IoT"
sizes = [60, 90]  # Percentage of access

plt.figure(figsize=(4, 4))
colors = ["#ff9999", "#66b3ff"]
explode = (0.05, 0.05)
plt.pie(
    sizes,
    labels=labels,
    autopct="%1.1f%%",
    startangle=90,
    colors=colors,
    explode=explode,
)
plt.title("Healthcare Access Increase", fontsize=14)
plt.tight_layout()
graph_path = "healthcare_access.png"
plt.savefig(graph_path, dpi=100)
plt.close()

# Add chart to slide
add_picture(
    slide, graph_path, left=Inches(3), top=Inches(2), width=Inches(5), height=Inches(4)
)

# Slide 14: Role of JRC Board
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Our Journey at JRC Board")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
points = ["• Key projects and innovations.", "• Collaborations with local communities."]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (Add images of team and products)
add_picture(
    slide,
    "jrc_board_projects.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(3),
    height=Inches(3),
)

# Slide 15: Role of FronTech LTD
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Innovations at FronTech LTD")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
points = ["• IoT solutions developed for sustainability.", "• Impact stories."]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (Add images of products and solutions)
add_picture(
    slide,
    "frontech_innovations.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(3),
    height=Inches(3),
)

# Slide 16: Encouraging Innovation
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Encouraging Innovation Among Students")

# Add background image (Placeholder)
add_picture(
    slide,
    "students_innovation.jpg",
    left=Inches(0),
    top=Inches(1.5),
    width=Inches(10),
    height=Inches(5),
)

# Slide 17: Global Exposure
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Global Opportunities")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "• Collaborations with international organizations.",
    "• Participating in global tech challenges.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (World map highlighting countries)
add_picture(
    slide,
    "world_map.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(4),
    height=Inches(3),
)

# Slide 18: Challenges and Solutions
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Challenges in Implementing IoT Solutions")

# Left Column: Challenges
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
p = left_text_frame.add_paragraph()
run = p.add_run()
run.text = "Challenges:"
run.font.name = "Calibri"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = TITLE_COLOR
challenges = [
    "• Infrastructure limitations.",
    "• High costs.",
    "• Need for skilled personnel.",
]
for challenge in challenges:
    p = left_text_frame.add_paragraph()
    p.text = challenge
    p.level = 1
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = CONTENT_COLOR

# Right Column: Solutions
right_text_frame = slide.placeholders[2].text_frame
right_text_frame.clear()
p = right_text_frame.add_paragraph()
run = p.add_run()
run.text = "Solutions:"
run.font.name = "Calibri"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = TITLE_COLOR
solutions = [
    "• Government support.",
    "• Local manufacturing.",
    "• Education and training.",
]
for solution in solutions:
    p = right_text_frame.add_paragraph()
    p.text = solution
    p.level = 1
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = CONTENT_COLOR

# Slide 19: Education and Skill Development
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Education and Skill Development")

content = [
    "• Importance of STEM education.",
    "• Training programs and workshops offered.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Add image of a workshop (Placeholder)
add_picture(
    slide,
    "workshop.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 20: Collaborations and Partnerships
slide = prs.slides.add_slide(two_content_layout)
add_custom_title(slide, "Building Collaborations")

# Left Content
left_text_frame = slide.placeholders[1].text_frame
left_text_frame.clear()
points = [
    "• Partnerships with universities and industry.",
    "• Benefits of collaborative innovation.",
]
for point in points:
    p = left_text_frame.add_paragraph()
    p.text = point
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Right Content (Logos of partners)
add_picture(
    slide,
    "partner_logos.jpg",
    left=Inches(6),
    top=Inches(1.5),
    width=Inches(3),
    height=Inches(3),
)

# Slide 21: Future Outlook
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "The Future of Robotics and IoT in Bangladesh")

content = [
    "• Emerging technologies (AI integration, 5G).",
    "• Vision for the next decade.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for line in content:
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = CONTENT_COLOR

# Add image of a futuristic smart city (Placeholder)
add_picture(
    slide,
    "future_smart_city.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 22: How You Can Get Involved
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "How You Can Get Involved")

content = [
    "Opportunities:",
    "• Internships.",
    "• Research projects.",
    "• Innovation challenges.",
    "Steps to participate.",
]

text_frame = slide.shapes.placeholders[1].text_frame
add_custom_content(text_frame, content)

# Add image of enthusiastic students (Placeholder)
add_picture(
    slide,
    "students_get_involved.jpg",
    left=Inches(6),
    top=Inches(2),
    width=Inches(3),
    height=Inches(3),
)

# Slide 23: Call to Action
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Join the Movement")

# Add inspirational quote
quote_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
quote_tf = quote_shape.text_frame
p = quote_tf.paragraphs[0]
p.text = '"Innovation is seeing what everybody has seen and thinking what nobody has thought." - Dr. Albert Szent-Györgyi'
p.font.name = "Calibri"
p.font.size = Pt(24)
p.font.italic = True
p.font.color.rgb = CONTENT_COLOR
p.alignment = PP_ALIGN.CENTER

# Add background image (Placeholder)
add_picture(
    slide,
    "call_to_action.jpg",
    left=Inches(0),
    top=Inches(4),
    width=Inches(10),
    height=Inches(2),
)

# Slide 24: Q&A Session
slide = prs.slides.add_slide(blank_slide_layout)
add_custom_title(slide, "Questions & Answers")

# Add background image (Placeholder)
add_picture(
    slide,
    "q_and_a.jpg",
    left=Inches(0),
    top=Inches(1.5),
    width=Inches(10),
    height=Inches(5),
)

# Slide 25: Thank You
slide = prs.slides.add_slide(title_and_content_layout)
add_custom_title(slide, "Thank You")

content = [
    "Contact Information:",
    "Email: your.email@example.com",
    "LinkedIn: linkedin.com/in/yourprofile",
    "Website: www.yourcompanywebsite.com",
    "Invitation to connect and collaborate.",
]

text_frame = slide.shapes.placeholders[1].text_frame
text_frame.clear()
for idx, line in enumerate(content):
    p = text_frame.add_paragraph()
    p.text = line
    p.font.name = "Calibri"
    if idx == 0:
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
    else:
        p.font.size = Pt(20)
        p.font.color.rgb = CONTENT_COLOR

# Add company logos or professional photo (Placeholder)
add_picture(
    slide, "logos.jpg", left=Inches(2), top=Inches(5), width=Inches(6), height=Inches(1)
)

# Save the presentation
prs.save("Empowering_Innovation.pptx")
